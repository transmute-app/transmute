import os
import subprocess  # nosec B404
import sys
import tempfile
import html as html_module
from pathlib import Path
from typing import Optional

from core import validate_safe_path

from .converter_interface import ConverterInterface


class LibreOfficeConverter(ConverterInterface):
    """
    Converter for presentation formats using LibreOffice.
    Uses LibreOffice's headless mode for server-side document conversion.

    Supports conversions between PowerPoint, ODP, PDF, image, and text formats.

    For PPTX input to text/HTML, uses python-pptx for cleaner extraction.
    For image outputs, routes through a PDF intermediary so all slides are
    included in the final image.
    """

    # Presentation formats that LibreOffice Impress can open
    supported_input_formats: set = {
        'pptx',
        'odp',
        'ppt',
        'pps',
        'ppsx',
        'pot',
        'potx',
        'pptm',
        'key',
    }

    # Formats that LibreOffice Impress can export to
    supported_output_formats: set = {
        'pptx',
        'pdf',
        'html',
        'txt',
        'odp',
        'ppt',
        'png',
        'jpeg',
    }
    qualities = {
        'low',
        'medium',
        'high',
    }

    # LibreOffice binary paths by platform
    _soffice_paths = {
        'darwin': '/Applications/LibreOffice.app/Contents/MacOS/soffice',
        'linux': '/usr/bin/soffice',
        'win32': 'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
    }
    soffice_path = _soffice_paths.get(sys.platform, 'soffice')

    # Map internal format names to LibreOffice --convert-to format strings.
    _lo_format_map = {
        'pptx': 'pptx',
        'ppt': 'ppt',
        'pdf': 'pdf',
        'html': 'html',
        'txt': 'txt',
        'odp': 'odp',
        'png': 'png',
        'jpeg': 'jpg',
        'eps': 'eps',
    }

    # Input formats that python-pptx can read natively (OOXML variants)
    _pptx_native_formats = {'pptx', 'pptm'}

    def __init__(self, input_file: str, output_dir: str, input_type: str, output_type: str):
        """
        Initialize LibreOffice converter.

        Args:
            input_file: Path to the input presentation file
            output_dir: Directory where the converted file will be saved
            input_type: Input file format (e.g., 'pptx', 'odp', 'ppt')
            output_type: Output file format (e.g., 'pdf', 'png', 'pptx')
        """
        super().__init__(input_file, output_dir, input_type, output_type)

    @classmethod
    def can_register(cls) -> bool:
        """
        Check if LibreOffice is available on the system.

        Returns:
            True if LibreOffice is installed and accessible, False otherwise.
        """
        try:
            subprocess.run(  # nosec B603
                [cls.soffice_path, '--headless', '--version'],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )
            return True
        except (subprocess.CalledProcessError, FileNotFoundError):
            return False

    def can_convert(self) -> bool:
        """
        Check if the input file can be converted to the output format.

        Returns:
            True if conversion is possible, False otherwise.
        """
        input_fmt = self.input_type.lower()
        output_fmt = self.output_type.lower()

        if input_fmt not in self.supported_input_formats:
            return False
        if output_fmt not in self.supported_output_formats:
            return False
        if input_fmt == output_fmt:
            return False

        return True

    @classmethod
    def get_formats_compatible_with(cls, format_type: str) -> set:
        """
        Get the set of compatible output formats for a given input format.

        Args:
            format_type: The input format to check compatibility for.

        Returns:
            Set of compatible output formats.
        """
        fmt = format_type.lower()
        if fmt not in cls.supported_input_formats:
            return set()
        return cls.supported_output_formats - {fmt}

    # ------------------------------------------------------------------
    # Main entry point
    # ------------------------------------------------------------------

    def convert(self, overwrite: bool = True, quality: Optional[str] = None) -> list[str]:
        """
        Convert the input presentation file.

        Routing:
        - PPTX/PPTM → TXT or HTML: python-pptx (clean text extraction)
        - Any → JPEG/PNG/EPS: PDF intermediary + PyMuPDF + Pillow (all pages)
        - Everything else: LibreOffice headless (direct conversion)

        Args:
            overwrite: Whether to overwrite existing output file (default: True)
            quality: Not applicable for presentation conversions, ignored.

        Returns:
            List containing the path to the converted output file.

        Raises:
            FileNotFoundError: If input file doesn't exist.
            ValueError: If the conversion is not supported.
            RuntimeError: If conversion fails.
        """
        if not self.can_convert():
            raise ValueError(
                f"Conversion from {self.input_type} to {self.output_type} "
                "is not supported."
            )

        if not os.path.isfile(self.input_file):
            raise FileNotFoundError(f"Input file not found: {self.input_file}")

        try:
            if self.output_type in ('txt', 'html'):
                if self.input_type in self._pptx_native_formats:
                    return self._convert_text_with_pptx(overwrite)
                return self._convert_text_via_pptx(overwrite)

            if self.output_type in ('jpeg', 'png', 'eps'):
                return self._convert_to_image(overwrite)

            return self._convert_with_libreoffice(overwrite)

        except subprocess.CalledProcessError as e:
            raise RuntimeError(
                f"LibreOffice conversion failed: {e.stderr or e.stdout or str(e)}"
            )
        except (ValueError, RuntimeError):
            raise
        except Exception as e:
            raise RuntimeError(f"Presentation conversion failed: {str(e)}")

    # ------------------------------------------------------------------
    # python-pptx text / HTML extraction (PPTX / PPTM only)
    # ------------------------------------------------------------------

    def _convert_text_with_pptx(self, overwrite: bool) -> list[str]:
        """Extract text or generate HTML from a PPTX file via python-pptx."""
        from pptx import Presentation

        input_filename = Path(self.input_file).stem
        output_file = os.path.join(
            self.output_dir, f"{input_filename}.{self.output_type}"
        )

        if not overwrite and os.path.exists(output_file):
            return [output_file]

        prs = Presentation(self.input_file)

        if self.output_type == 'txt':
            content = self._extract_text(prs)
        else:
            content = self._generate_html(prs)

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(content)

        return [output_file]

    def _convert_text_via_pptx(self, overwrite: bool) -> list[str]:
        """Convert non-PPTX presentations to TXT/HTML by going through PPTX first."""
        from pptx import Presentation

        input_filename = Path(self.input_file).stem
        output_file = os.path.join(
            self.output_dir, f"{input_filename}.{self.output_type}"
        )

        if not overwrite and os.path.exists(output_file):
            return [output_file]

        with tempfile.TemporaryDirectory() as tmp_dir:
            # Step 1: convert to PPTX via LibreOffice
            pptx_path = self._run_libreoffice(tmp_dir, 'pptx')

            # Step 2: extract text/HTML from the intermediate PPTX
            prs = Presentation(pptx_path)

        if self.output_type == 'txt':
            content = self._extract_text(prs)
        else:
            content = self._generate_html(prs)

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(content)

        return [output_file]

    @staticmethod
    def _extract_text(prs) -> str:
        """Return plain-text content for every slide."""
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = '\t'.join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            parts.append(row_text)
            if parts:
                slides_text.append(
                    f"--- Slide {i} ---\n" + '\n'.join(parts)
                )
        return '\n\n'.join(slides_text)

    @staticmethod
    def _generate_html(prs) -> str:
        """Return a clean, readable HTML document from all slides."""
        sections: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f'<section class="slide" id="slide-{i}">',
                     f'<h2>Slide {i}</h2>']
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(
                                f'<p>{html_module.escape(text)}</p>'
                            )
                if shape.has_table:
                    parts.append(
                        '<table border="1" cellpadding="4" '
                        'cellspacing="0">'
                    )
                    for ri, row in enumerate(shape.table.rows):
                        tag = 'th' if ri == 0 else 'td'
                        cells = ''.join(
                            f'<{tag}>{html_module.escape(cell.text.strip())}'
                            f'</{tag}>'
                            for cell in row.cells
                        )
                        parts.append(f'<tr>{cells}</tr>')
                    parts.append('</table>')
            parts.append('</section>')
            sections.append('\n'.join(parts))

        body = '\n<hr>\n'.join(sections)
        return (
            '<!DOCTYPE html>\n<html lang="en">\n<head>\n'
            '<meta charset="utf-8">\n'
            '<meta name="viewport" '
            'content="width=device-width, initial-scale=1.0">\n'
            '<title>Presentation</title>\n'
            '<style>\n'
            'body { font-family: sans-serif; max-width: 960px; '
            'margin: 0 auto; padding: 2rem; }\n'
            '.slide { margin: 2rem 0; padding: 1rem; }\n'
            'table { border-collapse: collapse; margin: 1rem 0; }\n'
            'th, td { padding: 0.5rem; text-align: left; '
            'border: 1px solid #ccc; }\n'
            'th { background-color: #f0f0f0; }\n'
            'hr { margin: 2rem 0; }\n'
            '</style>\n</head>\n'
            f'<body>\n{body}\n</body>\n</html>'
        )

    # ------------------------------------------------------------------
    # Image output via PDF intermediary (all input formats, all pages)
    # ------------------------------------------------------------------

    def _convert_to_image(self, overwrite: bool) -> list[str]:
        """Render every slide as an image by going PPTX → PDF → image."""
        import fitz  # PyMuPDF
        from PIL import Image

        input_filename = Path(self.input_file).stem
        output_file = os.path.join(
            self.output_dir, f"{input_filename}.{self.output_type}"
        )

        if not overwrite and os.path.exists(output_file):
            return [output_file]

        with tempfile.TemporaryDirectory() as tmp_dir:
            # Step 1: produce a PDF via LibreOffice
            pdf_path = self._run_libreoffice(tmp_dir, 'pdf')

            # Step 2: render every page at 200 DPI
            doc = fitz.open(pdf_path)
            page_images: list[Image.Image] = []
            zoom = 200 / 72
            matrix = fitz.Matrix(zoom, zoom)

            for page in doc:
                pix = page.get_pixmap(matrix=matrix)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                page_images.append(img)
            doc.close()

        if not page_images:
            raise RuntimeError("No pages were rendered from the presentation")

        # Step 3: stitch all pages into a single tall image
        total_width = max(img.width for img in page_images)
        total_height = sum(img.height for img in page_images)

        # libjpeg cannot handle dimensions > 65500 pixels.  Scale the
        # combined image down when the stitched height would exceed that.
        max_dim = 65500
        if total_height > max_dim or total_width > max_dim:
            scale = min(max_dim / total_height, max_dim / total_width)
            new_w = int(total_width * scale)
            new_h = int(total_height * scale)
            resized: list[Image.Image] = []
            for img in page_images:
                resized.append(
                    img.resize(
                        (int(img.width * scale), int(img.height * scale)),
                        Image.LANCZOS,
                    )
                )
            page_images = resized
            total_width, total_height = new_w, new_h

        combined = Image.new('RGB', (total_width, total_height), 'white')

        y_offset = 0
        for img in page_images:
            x_offset = (total_width - img.width) // 2
            combined.paste(img, (x_offset, y_offset))
            y_offset += img.height

        if self.output_type == 'jpeg':
            combined.save(output_file, 'JPEG', quality=95)
        elif self.output_type == 'png':
            combined.save(output_file, 'PNG')
        elif self.output_type == 'eps':
            combined.save(output_file, 'EPS')

        if not os.path.exists(output_file):
            raise RuntimeError(f"Output image was not created: {output_file}")

        return [output_file]

    # ------------------------------------------------------------------
    # Direct LibreOffice conversion (presentation ↔ presentation, PDF)
    # ------------------------------------------------------------------

    def _convert_with_libreoffice(self, overwrite: bool) -> list[str]:
        """Standard conversion via ``soffice --convert-to``."""
        lo_format = self._lo_format_map.get(self.output_type, self.output_type)
        input_filename = Path(self.input_file).stem
        output_file = os.path.join(
            self.output_dir, f"{input_filename}.{lo_format}"
        )

        if not overwrite and os.path.exists(output_file):
            return [output_file]

        validate_safe_path(self.input_file)
        validate_safe_path(output_file)

        self._run_libreoffice(self.output_dir, lo_format)

        if not os.path.exists(output_file):
            raise RuntimeError(f"Output file was not created: {output_file}")

        return [output_file]

    # ------------------------------------------------------------------
    # Shared LibreOffice invocation helper
    # ------------------------------------------------------------------

    def _run_libreoffice(self, output_dir: str, lo_format: str) -> str:
        """
        Run LibreOffice headless to convert ``self.input_file``.

        Returns:
            Path to the file produced by LibreOffice.
        """
        input_filename = Path(self.input_file).stem
        output_path = os.path.join(output_dir, f"{input_filename}.{lo_format}")

        with tempfile.TemporaryDirectory() as user_install_dir:
            cmd = [
                self.soffice_path,
                '--headless',
                '--norestore',
                f'-env:UserInstallation=file://{user_install_dir}',
                '--convert-to', lo_format,
                '--outdir', output_dir,
                self.input_file,
            ]

            subprocess.run(  # nosec B603
                cmd,
                capture_output=True,
                text=True,
                check=True,
                timeout=120,
            )

        if not os.path.exists(output_path):
            raise RuntimeError(
                f"LibreOffice did not produce the expected file: {output_path}"
            )

        return output_path
